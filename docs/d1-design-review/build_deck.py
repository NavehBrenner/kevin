"""Build the D1 Design Review .pptx deck for the AI-Assisted Teleoperation project.

This script renders a 17-slide deck using python-pptx, embedding the rasterized
architecture.svg and sequence.svg diagrams.

Setup (one-time):
    From the `code/` directory:
        uv pip install -e ".[docs]"     # if you've added the docs extra
    or quick:
        uv pip install python-pptx cairosvg

Run:
    cd code/docs/d1-design-review
    uv run python build_deck.py

Output:
    design-review.pptx (in this directory)

Notes:
    - If cairosvg fails to install (it has system deps on some boxes), pre-rasterize
      the SVGs manually with rsvg-convert or inkscape:
          rsvg-convert -w 2000 architecture.svg -o architecture.png
          rsvg-convert -w 2000 sequence.svg     -o sequence.png
      The script picks up the .png files automatically if they exist next to the .svgs.
"""

from __future__ import annotations

import io
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

try:
    import cairosvg  # type: ignore

    HAS_CAIROSVG = True
except ImportError:
    HAS_CAIROSVG = False


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

HERE = Path(__file__).parent
OUT_PATH = HERE / "design-review.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette (matches diagrams)
NAVY = RGBColor(0x1E, 0x27, 0x61)
CHARCOAL = RGBColor(0x21, 0x25, 0x29)
MUTED = RGBColor(0x6C, 0x75, 0x7D)
AMBER = RGBColor(0xD9, 0x77, 0x06)
GREEN = RGBColor(0x16, 0x65, 0x34)
INDIGO = RGBColor(0x37, 0x30, 0xA3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)

HEADER_FONT = "Calibri"
BODY_FONT = "Calibri"


def get_diagram_png(svg_path: Path, output_width_px: int = 2000) -> bytes:
    """Get PNG bytes for a diagram.

    Strategy: if a .png next to the .svg exists, use it. Otherwise rasterize via
    cairosvg. Otherwise raise with instructions.
    """
    png_path = svg_path.with_suffix(".png")
    if png_path.exists():
        return png_path.read_bytes()
    if HAS_CAIROSVG:
        return cairosvg.svg2png(url=str(svg_path), output_width=output_width_px)
    raise RuntimeError(
        f"Neither {png_path.name} exists nor is cairosvg installed.\n"
        f"Install cairosvg, OR pre-convert with:\n"
        f"  rsvg-convert -w {output_width_px} {svg_path.name} -o {png_path.name}"
    )


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_text(tf, text: str, *, size: int, bold: bool = False, color=CHARCOAL,
             font=BODY_FONT, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text(slide, text: str, *, left, top, width, height,
             size: int, bold=False, color=CHARCOAL, font=BODY_FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    set_text(tf, text, size=size, bold=bold, color=color, font=font, align=align)
    return box


def add_bullets(slide, items, *, left, top, width, height,
                size: int = 16, bold=False, color=CHARCOAL,
                bullet_char="•"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.name = BODY_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        p.space_after = Pt(6)
    return box


def add_title(slide, title: str, *, size: int = 32, color=NAVY):
    return add_text(
        slide, title,
        left=Inches(0.6), top=Inches(0.4),
        width=Inches(12), height=Inches(0.8),
        size=size, bold=True, color=color, font=HEADER_FONT,
    )


def add_subtitle(slide, subtitle: str, *, top=Inches(1.1)):
    return add_text(
        slide, subtitle,
        left=Inches(0.6), top=top,
        width=Inches(12), height=Inches(0.5),
        size=16, color=MUTED, font=HEADER_FONT,
    )


def add_image_centered(slide, image_bytes: bytes, *, top=Inches(1.6),
                       max_width=Inches(12.5), max_height=Inches(5.5)):
    bio = io.BytesIO(image_bytes)
    pic = slide.shapes.add_picture(bio, Inches(0), top, width=max_width)
    if pic.height > max_height:
        ratio = max_height / pic.height
        pic.width = int(pic.width * ratio)
        pic.height = max_height
    pic.left = int((SLIDE_W - pic.width) / 2)
    return pic


def add_notes(slide, notes_text: str):
    slide.notes_slide.notes_text_frame.text = notes_text


# ---------------------------------------------------------------------------
# Slide content
# ---------------------------------------------------------------------------

def slide_title(prs):
    s = blank_slide(prs)
    add_text(
        s, "AI-Assisted Robotic Teleoperation",
        left=Inches(0.8), top=Inches(2.4),
        width=Inches(11.7), height=Inches(1.4),
        size=46, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_text(
        s, "Precision peg-in-hole insertion with a vision-conditioned residual policy",
        left=Inches(0.8), top=Inches(3.5),
        width=Inches(11.7), height=Inches(0.8),
        size=22, color=CHARCOAL, font=BODY_FONT,
    )
    add_text(
        s, "Naveh Brenner   ·   OpenU Workshop in Autonomous Systems Simulation (20973)",
        left=Inches(0.8), top=Inches(5.8),
        width=Inches(11.7), height=Inches(0.4),
        size=15, color=MUTED, font=BODY_FONT,
    )
    add_text(
        s, "Fall 2026   ·   D1 Design Review",
        left=Inches(0.8), top=Inches(6.2),
        width=Inches(11.7), height=Inches(0.4),
        size=15, color=MUTED, font=BODY_FONT,
    )
    add_notes(
        s,
        "Opening slide.\n"
        "Hook: this is an AI-assisted teleoperation system - coarse human commands + AI sub-mm corrections - "
        "implemented for precision peg-in-hole insertion in MuJoCo.",
    )
    return s


def slide_problem(prs):
    s = blank_slide(prs)
    add_title(s, "The problem")
    add_subtitle(s, "Why neither human-only nor full-autonomy is satisfying for precision insertion")
    add_bullets(
        s,
        [
            "Precision insertion requires sub-millimeter accuracy at the end-effector.",
            "Human teleoperation alone is too coarse — hand tremor, latency, depth ambiguity, no haptic feedback.",
            "Full autonomy is brittle — perception edge cases, unseen geometries, no graceful degradation.",
            "Shared autonomy bridges the gap: the human supplies intent and goal context; the AI supplies the precision.",
        ],
        left=Inches(0.8), top=Inches(2.2), width=Inches(11.7), height=Inches(4.5),
        size=20,
    )
    add_notes(
        s,
        "Frame the gap. Natural impulse is 'just make the robot autonomous' or 'just hand it to a human'. "
        "Both fail for precision tasks. We're picking the middle ground deliberately.",
    )
    return s


def slide_what_we_build(prs):
    s = blank_slide(prs)
    add_title(s, "What we're building")
    add_subtitle(s, "A simulated Franka Panda performs peg-in-hole insertion under shared control")
    add_bullets(
        s,
        [
            "MuJoCo scene: Franka Panda arm + wall with multiple chamfered holes + pre-grasped peg.",
            "Wrist-mounted RGB camera + spotlight; 6-DoF wrist F/T sensor.",
            "Human supplies coarse 6-DoF gestures via webcam-tracked hand motion (MediaPipe).",
            "A behavioral-cloning residual policy adds sub-millimeter corrections in real time.",
            "All Python end-to-end — no C/C++/Rust extensions, no ROS.",
        ],
        left=Inches(0.8), top=Inches(2.2), width=Inches(11.7), height=Inches(4.5),
        size=20,
    )
    add_notes(
        s,
        "One sentence per bullet. The goal of this slide is to make the project concrete enough that the audience "
        "knows what's running on screen during the demo later.",
    )
    return s


def slide_goals(prs):
    s = blank_slide(prs)
    add_title(s, "Goals and requirements")
    add_bullets(
        s,
        [
            "Functional: coarse-from-human / fine-from-AI; quantifiably better than human-only.",
            "Architectural: clean strategy seams (input swappable; assistance swappable).",
            "Two-phase scope: Phase 1 F/T-only residual; Phase 2 vision-conditioned upgrade.",
            "Reliability: contact forces bounded by impedance + a force-cap watchdog — never crushes peg or wall.",
            "Reproducibility: deterministic evaluation; paired seeds across human-only / heuristic / learned modes.",
        ],
        left=Inches(0.8), top=Inches(2.0), width=Inches(11.7), height=Inches(4.8),
        size=20,
    )
    add_notes(
        s,
        "These derive from project-scope.md. The 'two-phase' point is the key scoping decision: "
        "F/T-only is the floor; vision is the upside.",
    )
    return s


def slide_kpis(prs):
    s = blank_slide(prs)
    add_title(s, "How we'll measure success")
    add_subtitle(s, "Three modes, same trials, paired seeds — measured by these KPIs")

    add_text(
        s, "Per-trial KPIs",
        left=Inches(0.8), top=Inches(2.2), width=Inches(5.5), height=Inches(0.5),
        size=22, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "Success rate (peg fully inserted within time + force budget)",
            "Time-to-insert (s)",
            "Peak contact force (N) — must stay under F_cap",
            "Contact count / contact-time fraction",
            "End-effector path smoothness (Δ-jerk)",
        ],
        left=Inches(0.8), top=Inches(2.8), width=Inches(5.8), height=Inches(4.0),
        size=17,
    )

    add_text(
        s, "Comparison setup",
        left=Inches(7.0), top=Inches(2.2), width=Inches(5.5), height=Inches(0.5),
        size=22, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "Three modes: human-only · heuristic · learned residual",
            "Paired seeds — same noisy-human trial under all three modes",
            "~100 trials per mode per configuration",
            "Target: F/T residual beats human-only on success rate AND peak force",
            "Eval harness is a passive observer (decoupled from controller)",
        ],
        left=Inches(7.0), top=Inches(2.8), width=Inches(5.8), height=Inches(4.0),
        size=17,
    )
    add_notes(
        s,
        "KPI target values are deliberately deferred to M6 (calibration milestone) - we want the task to be "
        "genuinely hard for human-only. The slide makes the WHAT clear; the numbers will land in M6.",
    )
    return s


def slide_high_level_approach(prs):
    s = blank_slide(prs)
    add_title(s, "High-level approach")
    add_subtitle(s, "Three layers, two swappable seams")

    items = [
        (
            "1.  Input layer",
            "Human-driven (webcam hand pose via MediaPipe, or keyboard fallback) or scripted noisy-human "
            "for training. All three behind one common interface.",
        ),
        (
            "2.  Assistance layer",
            "Human-only (no assist), heuristic (spiral search), or learned residual policy (behavioral cloning). "
            "All three behind one common interface — the swappable seam.",
        ),
        (
            "3.  Backbone controller",
            "Always-on. Differential IK + direction-dependent impedance + force-cap watchdog + lock-state machine. "
            "The safety backbone — never replaced.",
        ),
    ]

    y = Inches(2.0)
    for header, body in items:
        add_text(
            s, header,
            left=Inches(0.8), top=y, width=Inches(11.7), height=Inches(0.5),
            size=22, bold=True, color=NAVY, font=HEADER_FONT,
        )
        add_text(
            s, body,
            left=Inches(1.2), top=y + Inches(0.45),
            width=Inches(11.3), height=Inches(1.0),
            size=16, color=CHARCOAL, font=BODY_FONT,
        )
        y += Inches(1.55)

    add_notes(
        s,
        "Set up the architecture slide. 'Two swappable seams' is what the diagram on the next slide will show. "
        "Strategy-pattern at three places: input, assistance, lock (lock is internal to the controller).",
    )
    return s


def slide_architecture(prs):
    s = blank_slide(prs)
    add_title(s, "System architecture")
    img = get_diagram_png(HERE / "architecture.svg", output_width_px=2000)
    add_image_centered(s, img, top=Inches(1.2), max_width=Inches(12.5), max_height=Inches(6.0))
    add_notes(
        s,
        "Walk through:\n"
        "1. Input strategy on the left - three concrete implementations behind one interface.\n"
        "2. Assistance layer next - three concrete implementations, ditto.\n"
        "3. Controller and SimEnv are NOT swappable (backbone is always-on).\n"
        "4. Observation feedback closes the loop (dashed) back to assistance.\n"
        "5. Eval harness is a passive observer - no controller coupling.\n"
        "6. Bottom callout is the offline data-gen pipeline (separate from runtime).",
    )
    return s


def slide_sequence(prs):
    s = blank_slide(prs)
    add_title(s, "One control step")
    img = get_diagram_png(HERE / "sequence.svg", output_width_px=2000)
    add_image_centered(s, img, top=Inches(1.2), max_width=Inches(12.5), max_height=Inches(6.0))
    add_notes(
        s,
        "100 Hz control loop. Key safety property: the controller has its own force watchdog INDEPENDENT of the "
        "assistance layer's decisions - even a bad delta from the policy can't crush the peg.",
    )
    return s


def slide_ml_contribution(prs):
    s = blank_slide(prs)
    add_title(s, "What we actually learn")
    add_subtitle(s, "Residual policy trained by behavioral cloning on a privileged-info expert")

    add_text(
        s, "The model",
        left=Inches(0.8), top=Inches(2.0), width=Inches(5.5), height=Inches(0.5),
        size=22, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "Inputs: F/T history (~5 frames) + EE pose history + human command + (Phase 2) wrist RGB",
            "Output: clamped Δpose (Δx, Δy, Δz, ΔR) + Δgrip-force",
            "Architecture: small MLP (Phase 1) or frozen image backbone + MLP head (Phase 2)",
            "Runs at 100 Hz, real-time on CPU/GPU",
        ],
        left=Inches(0.8), top=Inches(2.6), width=Inches(5.9), height=Inches(3.5),
        size=15,
    )

    add_text(
        s, "Training data",
        left=Inches(7.0), top=Inches(2.0), width=Inches(5.5), height=Inches(0.5),
        size=22, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "Scripted noisy-human (programmatic, not a real person) drives the base command",
            "Analytical privileged-info expert produces the correction (asymmetric actor-critic pattern)",
            "Expert sees ground-truth peg + target hole pose; policy doesn't",
            "Coverage randomization per episode: hole positions, peg offset, noise pattern, "
            "(Phase 2) lighting/textures",
            "Target volume: hundreds of episodes; on-disk trajectory files",
        ],
        left=Inches(7.0), top=Inches(2.6), width=Inches(5.9), height=Inches(3.5),
        size=15,
    )

    add_text(
        s,
        "Why residual (not full replacement)?  Bounded deltas → safer; baseline always works; "
        "orthogonal failure modes.",
        left=Inches(0.8), top=Inches(6.4), width=Inches(11.7), height=Inches(0.6),
        size=14, color=MUTED, font=BODY_FONT, align=PP_ALIGN.LEFT,
    )
    add_notes(
        s,
        "Heart of the ML contribution. Asymmetric actor-critic (privileged-info expert teaches a sensor-limited "
        "policy) is well-established (e.g., Chen et al., 'Learning by Cheating'). We adapt it to a residual formulation.",
    )
    return s


def slide_two_phases(prs):
    s = blank_slide(prs)
    add_title(s, "Two phases, deliberately")
    add_subtitle(s, "Get a working baseline before betting on vision")

    add_text(
        s, "Phase 1 — F/T-only residual",
        left=Inches(0.8), top=Inches(2.0), width=Inches(5.7), height=Inches(0.5),
        size=22, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "Approximate hole prior + heuristic spiral search",
            "F/T residual corrects once peg is in contact",
            "No perception bandwidth needed — runs anywhere",
            "Goal: gate the architecture, get measured KPIs",
        ],
        left=Inches(0.8), top=Inches(2.6), width=Inches(5.7), height=Inches(3.5),
        size=16,
    )

    add_text(
        s, "Phase 2 — Vision-conditioned residual",
        left=Inches(7.0), top=Inches(2.0), width=Inches(5.7), height=Inches(0.5),
        size=22, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "Wrist RGB → policy can localize the target hole pre-contact",
            "Frozen pretrained image backbone + MLP head",
            "Same data-gen pipeline + image rendering on top",
            "Main ML contribution and the highest-risk milestone (M7)",
        ],
        left=Inches(7.0), top=Inches(2.6), width=Inches(5.7), height=Inches(3.5),
        size=16,
    )

    add_text(
        s,
        "Phase 1 is the floor — a polished Phase-1 project is a complete submission on its own. "
        "Phase 2 is the upside.",
        left=Inches(0.8), top=Inches(6.4), width=Inches(11.7), height=Inches(0.6),
        size=14, color=MUTED, font=BODY_FONT,
    )
    add_notes(
        s,
        "Risk-mitigation framing. The reviewer should walk away knowing we have a fall-back: even if vision "
        "doesn't pan out, Phase 1 is publishable.",
    )
    return s


def slide_alternatives(prs):
    s = blank_slide(prs)
    add_title(s, "Alternatives considered")
    add_subtitle(s, "Three design decisions, with the path-not-taken")

    headers = ["Decision", "Chosen", "Alternative considered", "Why chosen"]
    rows = [
        (
            "Policy form",
            "Residual (Δpose, Δgrip)",
            "Full-replacement policy",
            "Bounded deltas → safer; baseline always works; orthogonal failure modes.",
        ),
        (
            "Expert",
            "Analytical privileged-info",
            "RL with reward shaping",
            "Closed-form is faster + reproducible; sidesteps sim-to-real RL pain entirely.",
        ),
        (
            "Phase-1 sensing",
            "F/T only (+ heuristic prior)",
            "Direct to vision",
            "F/T-only ships fast and gates architecture; vision is the bigger bet, deferred to Phase 2.",
        ),
    ]

    table_left = Inches(0.6)
    table_top = Inches(2.0)
    col_widths = [Inches(2.3), Inches(3.0), Inches(3.0), Inches(4.4)]
    row_heights = [Inches(0.55), Inches(1.4), Inches(1.4), Inches(1.4)]

    rows_count = len(rows) + 1
    cols_count = len(headers)
    total_h = sum(row_heights, Inches(0))
    total_w = sum(col_widths, Inches(0))
    table_shape = s.shapes.add_table(rows_count, cols_count, table_left, table_top, total_w, total_h).table

    for i, w in enumerate(col_widths):
        table_shape.columns[i].width = w
    for i, h in enumerate(row_heights):
        table_shape.rows[i].height = h

    for i, header in enumerate(headers):
        cell = table_shape.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        tf = cell.text_frame
        tf.margin_left = Inches(0.1)
        tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.word_wrap = True
        set_text(tf, header, size=15, bold=True, color=WHITE, font=HEADER_FONT)

    for r, (decision, chosen, alt, why) in enumerate(rows, start=1):
        for c, val in enumerate([decision, chosen, alt, why]):
            cell = table_shape.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_BG
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.08)
            tf.margin_bottom = Inches(0.08)
            tf.word_wrap = True
            is_emphasis = (c == 1)
            set_text(
                tf, val,
                size=13,
                bold=(c == 0 or is_emphasis),
                color=NAVY if c == 0 else (AMBER if is_emphasis else CHARCOAL),
                font=BODY_FONT,
            )

    add_notes(
        s,
        "Three design choices, each defended. The booklet requires >=2 alternatives - we cover three across the "
        "most important axes: model form, expert form, and sensing.",
    )
    return s


def slide_scenarios(prs):
    s = blank_slide(prs)
    add_title(s, "Evaluation scenarios")
    add_subtitle(s, "Three difficulty tiers, controlled per-episode randomization")

    items = [
        ("Easy", AMBER, "Fixed hole location. Nominal lighting. Establishes the baseline upper bound."),
        (
            "Standard",
            NAVY,
            "3 distractor holes alongside target. Per-episode hole-offset noise (~5 mm). "
            "Initial peg offset (~10 mm). Nominal lighting. This is where headline KPIs are computed.",
        ),
        (
            "Hard (stretch)",
            INDIGO,
            "5+ distractors. Larger offset noise (~10 mm). Randomized lighting + wall textures (Phase 2 only). "
            "Probes robustness.",
        ),
    ]

    y = Inches(2.0)
    for name, color, body in items:
        add_text(
            s, name,
            left=Inches(0.8), top=y, width=Inches(2.8), height=Inches(0.5),
            size=22, bold=True, color=color, font=HEADER_FONT,
        )
        add_text(
            s, body,
            left=Inches(3.7), top=y, width=Inches(8.9), height=Inches(1.3),
            size=16, color=CHARCOAL, font=BODY_FONT, anchor=MSO_ANCHOR.TOP,
        )
        y += Inches(1.5)

    add_notes(
        s,
        "Per-episode randomizers: target hole index, hole position noise, initial peg offset, human-noise "
        "pattern, (Phase 2) lighting + textures. Magnitudes are deferred to M6 - calibrated so the task is "
        "genuinely hard for human-only.",
    )
    return s


def slide_eval_methodology(prs):
    s = blank_slide(prs)
    add_title(s, "Three-way KPI comparison")
    add_subtitle(s, "Paired seeds, passive observer, statistical analysis")

    add_bullets(
        s,
        [
            "Three modes per configuration: human-only · heuristic · learned residual.",
            "Paired seeds — each mode runs the same trial under the same noisy-human noise pattern.",
            "~100 trials per mode per configuration (300 trials per config total).",
            "Eval harness is a passive observer: subscribes to the observation stream, computes KPIs offline.",
            "Statistics: mean ± stderr per KPI per mode; paired Wilcoxon for pairwise mode comparison.",
            "Ablations (Phase 2): vision+F/T residual vs F/T-only residual on the same trials.",
        ],
        left=Inches(0.8), top=Inches(2.2), width=Inches(11.7), height=Inches(4.8),
        size=17,
    )

    add_notes(
        s,
        "Passive-observer pattern decouples eval from the controller - the controller doesn't know it's being "
        "evaluated. By design: lets us add new modes without changing the harness.",
    )
    return s


def slide_risks(prs):
    s = blank_slide(prs)
    add_title(s, "Risks and mitigations")

    risks = [
        (
            "R1 — Vision residual is data-hungry / overfits",
            "Mitigation: F/T baseline first (M5/M6) gives a known-working assist mode. Phase 2 adds domain "
            "randomization (lighting, textures) and ablates the vision contribution explicitly.",
        ),
        (
            "R2 — Sim-to-policy gap (state aliasing inside MuJoCo)",
            "Mitigation: aggressive coverage randomization in data generation — hole positions, peg offsets, "
            "noise patterns — so policy sees the test distribution at train time.",
        ),
        (
            "R3 — Scope vs. solo time budget (~225 h)",
            "Mitigation: explicit critical path M1→M6 (Phase 1 = floor); M7/M8 are upside. A polished "
            "Phase-1 project is a defensible submission on its own.",
        ),
        (
            "R4 — Live demo fragility (MediaPipe drift, lighting variance)",
            "Mitigation: keyboard fallback strategy; demo rehearses on canonical scene; demo is OPTIONAL — "
            "quantitative KPIs come from the scripted noisy-human, not the live demo.",
        ),
    ]

    y = Inches(2.0)
    for header, body in risks:
        add_text(
            s, header,
            left=Inches(0.8), top=y, width=Inches(11.7), height=Inches(0.5),
            size=17, bold=True, color=AMBER, font=HEADER_FONT,
        )
        add_text(
            s, body,
            left=Inches(1.2), top=y + Inches(0.4),
            width=Inches(11.3), height=Inches(0.8),
            size=14, color=CHARCOAL, font=BODY_FONT,
        )
        y += Inches(1.2)

    add_notes(
        s,
        "Four risks, ordered by impact. R3 (scope) is mitigated by the milestone roadmap itself - M6 is the "
        "latest point at which we can still pass even if M7 doesn't land.",
    )
    return s


def slide_timeline(prs):
    s = blank_slide(prs)
    add_title(s, "Milestone roadmap")
    add_subtitle(s, "~15 weeks remaining · ~10-15 hrs/week · M1 done")

    items = [
        ("late May", "M2 — Backbone controller", CHARCOAL, False),
        ("early June", "M3 — Heuristic assist online", CHARCOAL, False),
        ("mid June", "M4 — Expert + data generation (Foundation complete)", CHARCOAL, False),
        ("~Jun 8", "Topic approval (covered by project-scope.md)", MUTED, False),
        ("late June", "M5 — F/T-only residual policy", CHARCOAL, False),
        ("early July", "M6 — Eval harness + Phase 1 results (publishable)", GREEN, True),
        ("~mid July", "D1 Design Review — you are here", AMBER, True),
        ("mid-late July", "M7 — Vision-conditioned residual (Phase 2)", CHARCOAL, False),
        ("early August", "M8 — MediaPipe + keyboard input", CHARCOAL, False),
        ("mid-late August", "M9 — Final eval + polish", CHARCOAL, False),
        ("Aug 31", "D2 Final Submission", NAVY, True),
    ]

    y_top = Inches(1.9)
    row_h = Inches(0.42)
    for i, (date, label, color, bold) in enumerate(items):
        y = y_top + row_h * i
        add_text(
            s, date,
            left=Inches(0.8), top=y, width=Inches(2.5), height=row_h,
            size=14, bold=False, color=MUTED, font=BODY_FONT,
        )
        add_text(
            s, label,
            left=Inches(3.4), top=y, width=Inches(9.3), height=row_h,
            size=15, bold=bold, color=color, font=BODY_FONT,
        )

    add_notes(
        s,
        "Critical path M1->M6. D1 lands right after M6 so we should be presenting actual Phase-1 KPIs. "
        "If M5/M6 slip, D1 still shows M1-M4 working (heuristic-assist demo + data-gen pipeline).",
    )
    return s


def slide_current_state(prs):
    s = blank_slide(prs)
    add_title(s, "Where we are today")
    add_subtitle(s, "M1 complete — sim environment online behind a clean SimEnv API")

    add_text(
        s, "M1 Done",
        left=Inches(0.8), top=Inches(2.0), width=Inches(5.7), height=Inches(0.5),
        size=22, bold=True, color=GREEN, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "MuJoCo scene: Franka Panda + chamfered wall + pre-grasped peg",
            "Wrist RGB camera + wrist spotlight + 6-DoF F/T sensor",
            "Viewer (interactive) and headless render paths",
            "SimEnv class API + structured Observation dataclass",
            "Smoke tests passing; documented in milestone-1-spec.md",
        ],
        left=Inches(0.8), top=Inches(2.6), width=Inches(5.7), height=Inches(3.5),
        size=15,
    )

    add_text(
        s, "Next — M2",
        left=Inches(7.0), top=Inches(2.0), width=Inches(5.7), height=Inches(0.5),
        size=22, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "Backbone controller: differential IK + direction-dependent impedance",
            "Force-cap watchdog + lock-state machine (hold-lock / park-lock)",
            "Manual pose dev harness — drive the arm with keyboard targets",
            "Acceptance: bounded contact forces; smooth pose tracking",
        ],
        left=Inches(7.0), top=Inches(2.6), width=Inches(5.7), height=Inches(3.5),
        size=15,
    )

    add_text(
        s, "Demo plan (optional): live MuJoCo viewer showing the scene + SimEnv API walkthrough.",
        left=Inches(0.8), top=Inches(6.4), width=Inches(11.7), height=Inches(0.5),
        size=14, color=MUTED, font=BODY_FONT,
    )
    add_notes(
        s,
        "'Where are we' snapshot. By D1 (mid-July) we should have M2-M5 done, so this slide will be updated "
        "before the actual meeting. For now reflects end-of-May state.",
    )
    return s


def slide_qa(prs):
    s = blank_slide(prs)
    add_text(
        s, "Questions & discussion",
        left=Inches(0.8), top=Inches(1.5),
        width=Inches(11.7), height=Inches(1.0),
        size=42, bold=True, color=NAVY, font=HEADER_FONT,
    )
    add_text(
        s, "Specific things I'd appreciate feedback on:",
        left=Inches(0.8), top=Inches(3.0),
        width=Inches(11.7), height=Inches(0.5),
        size=18, bold=True, color=CHARCOAL, font=HEADER_FONT,
    )
    add_bullets(
        s,
        [
            "Architecture: does the assistance-layer abstraction support the comparisons you'd want to see?",
            "ML: any concerns about the privileged-info expert (asymmetric actor-critic) approach?",
            "Evaluation: preferred statistical tests, sample-size targets, additional KPIs?",
            "Phase split: is the F/T-only Phase 1 the right floor, or should Phase 1 already include vision?",
        ],
        left=Inches(0.8), top=Inches(3.6), width=Inches(11.7), height=Inches(3.5),
        size=16,
    )
    add_text(
        s, "Naveh Brenner · navegerc@gmail.com",
        left=Inches(0.8), top=Inches(6.6), width=Inches(11.7), height=Inches(0.4),
        size=12, color=MUTED, font=BODY_FONT,
    )
    add_notes(s, "End with explicit discussion prompts so reviewers know where you want their input.")
    return s


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_what_we_build(prs)
    slide_goals(prs)
    slide_kpis(prs)
    slide_high_level_approach(prs)
    slide_architecture(prs)
    slide_sequence(prs)
    slide_ml_contribution(prs)
    slide_two_phases(prs)
    slide_alternatives(prs)
    slide_scenarios(prs)
    slide_eval_methodology(prs)
    slide_risks(prs)
    slide_timeline(prs)
    slide_current_state(prs)
    slide_qa(prs)

    prs.save(str(OUT_PATH))
    print(f"Wrote: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
